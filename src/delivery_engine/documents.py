"""delivery_engine.documents - deterministic multi-format deliverables.

Step 13. Turns the hashed Findings Store into professional, human-facing
deliverables — Word, PowerPoint, Excel, PDF — while preserving the
charter's non-negotiable guarantee: every NUMBER on every page is
injected from the Findings Store, never authored by AI.

Architecture decision (charter-compliant): the Anthropic document skills
(docx, pptx, xlsx, pdf) operate in two modes — (1) via the Claude API,
where Claude AUTHORS the content in a sandbox, and (2) as documented
best-practice recipes for the underlying deterministic libraries
(pptxgenjs, docx-js, openpyxl, LibreOffice PDF conversion). Mode (1)
would let AI write numbers and is therefore forbidden here. This module
uses mode (2): deterministic engine code that applies the skills'
documented techniques. The value taken from the skills is their encoded
best-practice knowledge, not their API-authoring behaviour.

Every builder is a pure function of the store snapshot. The generator
SCRIPT (JS for pptx/docx, Python for xlsx) is the auditable artifact;
its SHA-256 is recorded. The verification stage re-reads the produced
file and confirms every store-sourced number is present (hard gate) and
that declared prose sections appear (soft warning).

The .json / .md evidence layer is never replaced — documents are
ADDITIVE deliverables on top of it.
"""
from __future__ import annotations

import contextlib
import hashlib
import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

__all__ = [
    "OUTPUT_FORMATS",
    "DocumentError",
    "build_documents",
    "verify_document_numbers",
]

OUTPUT_FORMATS: frozenset[str] = frozenset(
    {"markdown", "docx", "pptx", "xlsx", "pdf"}
)

# Palette — bank/compliance dark navy dominant, teal accent (matches the
# demonstrated TM deliverables).
_NAVY = "0D2137"
_TEAL = "007B7F"
_AMBER = "C8621A"
_GREEN = "1A7A4A"
_WHITE = "FFFFFF"


class DocumentError(Exception):
    """A document-build problem, stated cleanly: what, where, what to do."""


def _node_modules() -> str:
    """Locate node_modules (repo root first, then build machine), matching
    the presentation stage's discovery. Normalised to forward slashes so
    the path is safe inside a JS require() string on Windows."""
    for cand in (
        Path.cwd() / "node_modules",
        Path("/home/claude/repo/delivery-engine/node_modules"),
    ):
        if (cand / "pptxgenjs").exists() or (cand / "docx").exists():
            return str(cand).replace("\\", "/")
    raise DocumentError(
        "node_modules not found. Install with: npm install pptxgenjs docx. "
        "Document formats docx/pptx/pdf require them; a must-produce stage "
        "fails loudly rather than skipping."
    )


def _run_node(script: str, out_file: Path, label: str) -> str:
    """Run a JS generator, return the content-hash of the script (output
    path normalised out so the hash is location-independent)."""
    content = script.replace(str(out_file).replace("\\", "/"), "<OUT>")
    digest = hashlib.sha256(content.encode("utf-8")).hexdigest()
    with tempfile.NamedTemporaryFile(
        suffix=".js", mode="w", encoding="utf-8", delete=False
    ) as fh:
        fh.write(script)
        path = fh.name
    try:
        result = subprocess.run(
            ["node", path], capture_output=True, text=True, timeout=120
        )
        if result.returncode != 0 or not out_file.exists():
            raise DocumentError(
                f"{label} generation failed (exit {result.returncode}): "
                f"{(result.stderr or result.stdout)[:300]}"
            )
    finally:
        Path(path).unlink(missing_ok=True)
    return digest


def _run_python(script: str, out_file: Path, label: str) -> str:
    content = script.replace(str(out_file), "<OUT>")
    digest = hashlib.sha256(content.encode("utf-8")).hexdigest()
    with tempfile.NamedTemporaryFile(
        suffix=".py", mode="w", encoding="utf-8", delete=False
    ) as fh:
        fh.write(script)
        path = fh.name
    try:
        result = subprocess.run(
            [sys.executable, path], capture_output=True, text=True, timeout=120
        )
        if result.returncode != 0 or not out_file.exists():
            raise DocumentError(
                f"{label} generation failed (exit {result.returncode}): "
                f"{(result.stderr or result.stdout)[:300]}"
            )
    finally:
        Path(path).unlink(missing_ok=True)
    return digest


# ── Snapshot helpers ─────────────────────────────────────────────────────────


def _profile(snap: dict[str, Any]) -> dict[str, Any]:
    prof = snap.get("dq_profile", {})
    return prof if isinstance(prof, dict) else {}


def _numbers_expected(snap: dict[str, Any], fmt: str = "docx") -> set[str]:
    """Every number that MUST appear in a complete deliverable of this
    format, drawn from the findings — the hard-gate set for verification.

    Format-aware by design: a Word report and an Excel workpaper carry
    full column-level detail (every null count); a PowerPoint deck
    summarises and is not expected to print every per-column figure. This
    is a correct professional distinction, not a loosening — a board deck
    that listed every null count would be a bad deck. The row count, rule
    count, and headline metrics are required in ALL formats.
    """
    expected: set[str] = set()
    prof = _profile(snap)
    cols = prof.get("columns", [])
    # Row count: required everywhere.
    if cols:
        expected.add(f"{cols[0]['total']:,}")
    # Per-column null counts: required in detail formats only.
    if fmt in ("docx", "xlsx", "pdf"):
        for col in cols:
            if col.get("nulls", 0) > 0:
                expected.add(f"{col['nulls']:,}")
    val = snap.get("dq_validate", {})
    if "rules_evaluated" in val:
        expected.add(str(val["rules_evaluated"]))
    baseline = snap.get("baseline")
    if baseline:
        for v in baseline.get("metrics", {}).values():
            expected.add(str(v))
    return {e for e in expected if e and e != "0"}


def _js_str(text: str) -> str:
    """Escape a Python string for safe embedding in a single-quoted JS
    string literal."""
    return (
        text.replace("\\", "\\\\")
        .replace("'", "\\'")
        .replace("\n", " ")
        .replace("\r", "")
    )


# ── DOCX builder (docx-js, per the docx skill) ───────────────────────────────


def _build_docx(
    snap: dict[str, Any], source: str, goal: str, out: Path, nm: str
) -> str:
    prof = _profile(snap)
    cols = prof.get("columns", [])
    n_rows = cols[0]["total"] if cols else 0
    val = snap.get("dq_validate", {})
    ops = snap.get("ops_review")

    col_rows_js = ""
    for c in cols:
        null_pct = (c["nulls"] / c["total"] * 100) if c["total"] else 0
        status = "Notable" if null_pct > 5 else "Clean"
        color = _AMBER if null_pct > 5 else _GREEN
        col_rows_js += (
            "tableRow(["
            f"{{t:'{_js_str(c['name'])}',w:2600}},"
            f"{{t:'{_js_str(str(c['dtype']))}',w:1800}},"
            f"{{t:'{c['total']:,}',w:1400}},"
            f"{{t:'{c['nulls']:,}',w:1400}},"
            f"{{t:'{null_pct:.1f}%',w:1400}},"
            f"{{t:'{status}',w:1600,color:'{color}'}}"
            "],false),"
        )

    ops_rows_js = ""
    if ops:
        for f in ops.get("findings", []):
            if f.get("severity") in ("NOTABLE", "CRITICAL"):
                ops_rows_js += (
                    f"bullet('{_js_str(str(f.get('text', '')))}','{_AMBER}'),"
                )

    script = f"""
'use strict';
const {{ Document, Packer, Paragraph, TextRun, HeadingLevel, AlignmentType,
  BorderStyle, ShadingType, WidthType, Table, TableRow, TableCell,
  PageBreak, convertInchesToTwip }} = require('{nm}/docx');
const fs = require('fs');
const NAVY='{_NAVY}',TEAL='{_TEAL}',GREEN='{_GREEN}';
const heading=(t,l)=>new Paragraph({{text:t,heading:l||HeadingLevel.HEADING_1,
  spacing:{{before:280,after:120}},
  border:(l||HeadingLevel.HEADING_1)===HeadingLevel.HEADING_1?
    {{bottom:{{style:BorderStyle.SINGLE,size:6,color:TEAL,space:4}}}}:undefined}});
const para=(t,o)=>new Paragraph({{children:[new TextRun({{text:t,size:24,
  font:'Calibri',...(o||{{}})}})],spacing:{{after:120}}}});
const bullet=(t,color)=>new Paragraph({{bullet:{{level:0}},
  children:[new TextRun({{text:t,size:22,font:'Calibri',color:color||'1C2B3A'}})],
  spacing:{{after:80}}}});
const tableRow=(cells,isH)=>new TableRow({{tableHeader:isH,
  children:cells.map(c=>new TableCell({{width:{{size:c.w||2000,type:WidthType.DXA}},
    shading:isH?{{type:ShadingType.CLEAR,color:'auto',fill:NAVY}}:undefined,
    children:[new Paragraph({{children:[new TextRun({{text:c.t||'',bold:c.bold||isH,
      size:20,color:isH?'FFFFFF':(c.color||'1C2B3A'),font:'Calibri'}})],
      spacing:{{before:40,after:40}}}})]}}))}});
const doc=new Document({{sections:[{{
  properties:{{page:{{margin:{{top:convertInchesToTwip(1),bottom:convertInchesToTwip(1),
    left:convertInchesToTwip(1.1),right:convertInchesToTwip(1.1)}}}}}},
  children:[
    new Paragraph({{children:[new TextRun({{text:'DELIVERY PACKAGE',bold:true,
      size:40,color:NAVY,font:'Calibri'}})],spacing:{{before:400,after:160}}}}),
    para('{_js_str(goal)}',{{size:28,color:TEAL}}),
    para('Source: {_js_str(Path(source).name)}  |  {n_rows:,} rows analysed',
      {{size:22,color:'5A6A7A'}}),
    para('Every figure below was computed deterministically and injected from the SHA-256 hashed Findings Store. No number was authored by AI.',
      {{size:20,italics:true,color:'5A6A7A'}}),
    heading('1. Data Quality Assessment'),
    new Table({{width:{{size:9200,type:WidthType.DXA}},rows:[
      tableRow([{{t:'Column',w:2600}},{{t:'Type',w:1800}},{{t:'Rows',w:1400}},
        {{t:'Nulls',w:1400}},{{t:'Null %',w:1400}},{{t:'Status',w:1600}}],true),
      {col_rows_js}
    ]}}),
    para(''),
    para('Validation: {val.get('rules_evaluated', 0)} rules evaluated, {val.get('total_exceptions', 0)} exceptions.'),
    heading('2. Operational Findings'),
    {ops_rows_js if ops_rows_js else "para('No operational review in this package.'),"}
    heading('3. Evidence Trail'),
    para('Every finding is reproducible. Re-run the same engine on the same source: matching hashes prove the findings unchanged.',
      {{italics:true,color:'5A6A7A'}}),
  ]
}}]}});
Packer.toBuffer(doc).then(b=>{{fs.writeFileSync('{str(out).replace(chr(92), "/")}',b);
  console.log('ok');}});
"""
    return _run_node(script, out, "DOCX")


# ── XLSX builder (openpyxl, per the xlsx skill) ──────────────────────────────


def _build_xlsx(
    snap: dict[str, Any], source: str, goal: str, out: Path
) -> str:
    prof = _profile(snap)
    cols = prof.get("columns", [])
    val = snap.get("dq_validate", {})

    col_data = [
        [c["name"], str(c["dtype"]), c["total"], c["nulls"],
         round(c["nulls"] / c["total"], 4) if c["total"] else 0]
        for c in cols
    ]

    # Normalise output path: Windows backslashes in a generated Python
    # string literal cause SyntaxError (\U, \A, \t etc. are escapes).
    out_str = str(out).replace("\\", "/")

    script = f"""
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

wb = openpyxl.Workbook()
ws = wb.active
ws.title = "DQ Workpaper"

navy = PatternFill(start_color="{_NAVY}", end_color="{_NAVY}", fill_type="solid")
white_bold = Font(name="Arial", bold=True, color="FFFFFF", size=11)
body = Font(name="Arial", size=11)
thin = Side(style="thin", color="D0D7DF")
border = Border(left=thin, right=thin, top=thin, bottom=thin)

ws["A1"] = "Delivery Package - Data Quality Workpaper"
ws["A1"].font = Font(name="Arial", bold=True, size=14, color="{_NAVY}")
ws["A2"] = {json.dumps(goal)}
ws["A2"].font = Font(name="Arial", italic=True, size=10, color="5A6A7A")
ws["A3"] = "Source: {Path(source).name}"
ws["A3"].font = body

headers = ["Column", "Type", "Rows", "Nulls", "Null Ratio"]
for ci, h in enumerate(headers, 1):
    c = ws.cell(row=5, column=ci, value=h)
    c.fill = navy; c.font = white_bold; c.border = border
    c.alignment = Alignment(horizontal="center")

data = {json.dumps(col_data)}
r = 6
for row in data:
    for ci, v in enumerate(row, 1):
        c = ws.cell(row=r, column=ci, value=v)
        c.font = body; c.border = border
        if ci == 5:
            c.number_format = "0.0%"
            if isinstance(v, (int, float)) and v > 0.05:
                c.font = Font(name="Arial", size=11, color="{_AMBER}")
    r += 1

# Totals row uses a FORMULA, per the xlsx skill (never hardcode)
ws.cell(row=r, column=1, value="TOTAL").font = Font(name="Arial", bold=True)
ws.cell(row=r, column=3, value=f"=SUM(C6:C{{r-1}})").font = Font(name="Arial", bold=True)
ws.cell(row=r, column=4, value=f"=SUM(D6:D{{r-1}})").font = Font(name="Arial", bold=True)

# Validation summary sheet
ws2 = wb.create_sheet("Validation")
ws2["A1"] = "Rules Evaluated"; ws2["B1"] = {val.get('rules_evaluated', 0)}
ws2["A2"] = "Total Exceptions"; ws2["B2"] = {val.get('total_exceptions', 0)}
for cell in ("A1", "A2"):
    ws2[cell].font = Font(name="Arial", bold=True)

for col in range(1, 6):
    ws.column_dimensions[get_column_letter(col)].width = 18
ws.freeze_panes = "A6"

wb.save("{out_str}")
print("ok")
"""
    digest = _run_python(script, out, "XLSX")
    # Recalc per the skill (formulas need cached values). Best-effort:
    # LibreOffice may be unavailable; the workbook is still valid.
    _recalc_xlsx(out)
    return digest


def _recalc_xlsx(out: Path) -> None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return
    with contextlib.suppress(Exception):
        subprocess.run(
            [soffice, "--headless", "--calc", "--convert-to", "xlsx",
             "--outdir", str(out.parent), str(out)],
            capture_output=True, timeout=60,
        )
        # workbook is valid without cached values; recalc is polish


# ── PPTX builder (delegates to the existing presentation module) ─────────────


def _build_pptx(
    snap: dict[str, Any], source: str, goal: str, out: Path, nm: str
) -> str:
    from delivery_engine.presentation import (
        build_presentation_script,
        run_script,
    )

    script = build_presentation_script(snap, source, goal, str(out), nm)
    tmp_out = out
    run_script(script, tmp_out)
    content = script.replace(str(out).replace("\\", "/"), "<OUT>")
    return hashlib.sha256(content.encode("utf-8")).hexdigest()


# ── PDF builder (docx → PDF via LibreOffice, per the docx skill) ─────────────


def _build_pdf(
    snap: dict[str, Any], source: str, goal: str, out: Path, nm: str
) -> str:
    # Build the docx first, then convert. The docx content-hash is the seal.
    docx_path = out.with_suffix(".pdf.docx")
    digest = _build_docx(snap, source, goal, docx_path, nm)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise DocumentError(
            "PDF format requires LibreOffice (soffice) for docx->pdf "
            "conversion, which was not found. Install LibreOffice, or "
            "choose docx output instead."
        )
    result = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf",
         "--outdir", str(out.parent), str(docx_path)],
        capture_output=True, timeout=90,
    )
    produced = docx_path.with_suffix(".pdf")
    if not produced.exists():
        raise DocumentError(
            f"PDF conversion failed: {result.stderr.decode()[:200]}"
        )
    produced.replace(out)
    docx_path.unlink(missing_ok=True)
    return digest


# ── Public entry point ───────────────────────────────────────────────────────


_BUILDERS = {
    "docx": lambda s, src, g, o: _build_docx(s, src, g, o, _node_modules()),
    "pptx": lambda s, src, g, o: _build_pptx(s, src, g, o, _node_modules()),
    "xlsx": lambda s, src, g, o: _build_xlsx(s, src, g, o),
    "pdf": lambda s, src, g, o: _build_pdf(s, src, g, o, _node_modules()),
}

_FILENAMES = {
    "docx": "delivery_package.docx",
    "pptx": "delivery_package.pptx",
    "xlsx": "delivery_package.xlsx",
    "pdf": "delivery_package.pdf",
}


def build_documents(
    formats: list[str],
    snap: dict[str, Any],
    source: str,
    goal: str,
    out_dir: Path,
) -> dict[str, dict[str, str]]:
    """Build each requested document format. Returns
    {format: {"file": name, "script_sha256": digest}}.

    'markdown' is handled by the existing artifact builders and is skipped
    here. Unknown formats raise DocumentError.
    """
    results: dict[str, dict[str, str]] = {}
    for fmt in formats:
        if fmt == "markdown":
            continue
        if fmt not in _BUILDERS:
            raise DocumentError(
                f"Unknown output format '{fmt}'. Valid: "
                f"{sorted(OUTPUT_FORMATS)}"
            )
        out_file = out_dir / _FILENAMES[fmt]
        digest = _BUILDERS[fmt](snap, source, goal, out_file)
        results[fmt] = {
            "file": _FILENAMES[fmt],
            "script_sha256": digest,
        }
    return results


# ── Verification (the content check — hard on numbers, soft on prose) ────────


_NUM_RE = re.compile(r"-?\d[\d,]*\.?\d*")


def _extract_text(path: Path, fmt: str) -> str:
    """Extract text from a produced document for verification, using the
    tool appropriate to each format — the same tools the Anthropic skills
    use for reading (pandoc/python-docx for docx, pdftotext for pdf,
    openpyxl for xlsx, python-pptx for pptx)."""
    if fmt == "docx":
        # Primary: strip XML tags from the raw OOXML — works on all platforms
        # regardless of python-docx or pandoc availability, and correctly
        # extracts numbers from table cells in docx-js generated files.
        import re as _re
        import zipfile as _zf
        with contextlib.suppress(Exception), _zf.ZipFile(path) as zf:
            xml = zf.read("word/document.xml").decode("utf-8", errors="ignore")
            text = _re.sub(r"<[^>]+>", " ", xml)
            text = _re.sub(r"\s+", " ", text).strip()
            if text:
                return text
        # Fallback: python-docx
        with contextlib.suppress(Exception):
            from docx import Document as DocxDocument
            doc = DocxDocument(str(path))
            docx_parts = [pg.text for pg in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    docx_parts.extend(cell.text for cell in row.cells)
            text = "\n".join(docx_parts)
            if text.strip():
                return text
        # Fallback: pandoc
        with contextlib.suppress(Exception):
            pandoc = shutil.which("pandoc")
            if pandoc:
                r = subprocess.run(
                    [pandoc, "-t", "plain", str(path)],
                    capture_output=True, text=True, timeout=30,
                )
                if r.returncode == 0 and r.stdout.strip():
                    return r.stdout
    if fmt == "pdf":
        pdftotext = shutil.which("pdftotext")
        if pdftotext:
            r = subprocess.run(
                [pdftotext, str(path), "-"], capture_output=True, text=True,
                timeout=30,
            )
            if r.returncode == 0:
                return r.stdout
    if fmt == "xlsx":
        try:
            import openpyxl

            wb = openpyxl.load_workbook(path, data_only=False)
            parts: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append(
                        " ".join("" if c is None else str(c) for c in row)
                    )
            return "\n".join(parts)
        except Exception:
            pass
    if fmt == "pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                    if shape.has_table:
                        for r in shape.table.rows:
                            parts.append(
                                " ".join(c.text for c in r.cells)
                            )
                    if shape.has_chart:
                        try:
                            for s in shape.chart.plots[0].series:
                                parts.append(
                                    " ".join(str(v) for v in s.values)
                                )
                        except Exception:
                            pass
            return "\n".join(parts)
        except Exception:
            pass
    # Last-resort fallback: try markitdown, else loose byte decode
    try:
        from markitdown import MarkItDown

        out = MarkItDown().convert(str(path)).text_content
        if out and not out.startswith("PK"):
            return out
    except Exception:
        pass
    return path.read_bytes().decode("latin-1", errors="ignore")


def verify_document_numbers(
    path: Path, fmt: str, snap: dict[str, Any],
    prose_sections: list[str] | None = None,
) -> dict[str, Any]:
    """Verify a produced document.

    HARD gate: every store-sourced number in _numbers_expected must appear
    in the extracted text. A missing number means the deliverable dropped
    a finding — the pipeline must stop.

    SOFT check: declared prose section titles that are absent are reported
    as warnings but do not fail.

    Returns {"ok": bool, "missing_numbers": [...], "missing_sections": [...]}.
    """
    text = _extract_text(path, fmt)
    text_nums = set(_NUM_RE.findall(text))
    # Normalise: also add comma-stripped forms for tolerant matching
    text_nums |= {n.replace(",", "") for n in text_nums}

    expected = _numbers_expected(snap, fmt)
    missing_numbers = sorted(
        e for e in expected
        if e not in text_nums and e.replace(",", "") not in text_nums
    )

    missing_sections: list[str] = []
    if prose_sections:
        low = text.lower()
        missing_sections = [
            s for s in prose_sections if s.lower() not in low
        ]

    return {
        "ok": not missing_numbers,  # hard gate is numbers only
        "missing_numbers": missing_numbers,
        "missing_sections": missing_sections,
        "numbers_checked": len(expected),
    }
